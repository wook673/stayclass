from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# Colors
C_BG     = RGBColor(0x0A, 0x0A, 0x0A)
C_SURF   = RGBColor(0x14, 0x14, 0x14)
C_TEXT   = RGBColor(0xF0, 0xEC, 0xE4)
C_MUTED  = RGBColor(0x77, 0x77, 0x77)
C_LIME   = RGBColor(0xC8, 0xFF, 0x57)
C_BLUE   = RGBColor(0x38, 0xBD, 0xF8)
C_BORDER = RGBColor(0x28, 0x28, 0x28)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # completely blank


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color=C_BG):
    bg_shape = slide.shapes.add_shape(1, 0, 0, W, H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()


def rect(slide, x, y, w, h, color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, text, x, y, w, h,
          size=Pt(24), bold=False, color=C_TEXT,
          align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Noto Sans KR" if False else "맑은 고딕"
    return tb


def label(slide, text, x, y, w=Inches(10), size=Pt(13), color=C_LIME):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"


def title_h(slide, text, x, y, w=Inches(11), size=Pt(52), color=C_TEXT):
    tb = slide.shapes.add_textbox(x, y, w, Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = True
    run.font.color.rgb = color
    run.font.name = "맑은 고딕"
    return tb


def bullet_list(slide, items, x, y, w, size=Pt(20), color=C_MUTED, gap=Inches(0.45)):
    for i, item in enumerate(items):
        tb = slide.shapes.add_textbox(x, y + i * gap, w, Inches(0.5))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"— {item}"
        run.font.size = size
        run.font.color.rgb = color
        run.font.name = "맑은 고딕"


def stat_block(slide, num, unit, desc, x, y):
    txbox(slide, num, x, y, Inches(2), Inches(1.0), size=Pt(52), bold=True, color=C_LIME)
    txbox(slide, unit, x + Inches(1.5), y + Inches(0.1), Inches(1.2), Inches(0.6), size=Pt(22), color=C_BLUE)
    txbox(slide, desc, x, y + Inches(1.0), Inches(2.5), Inches(0.5), size=Pt(14), color=C_MUTED)


def divider(slide, y, color=C_BORDER):
    rect(slide, Inches(1), y, Inches(11.33), Pt(1), color)


# ─────────────────────────────────────────────
# SLIDE 1 — 표지
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_LIME)  # left accent bar

txbox(s, "STAYCLASS", Inches(1), Inches(0.8), Inches(6), Inches(0.5),
      size=Pt(14), bold=True, color=C_MUTED)

title_h(s, "직장 다니면서\n월 150만원\n더 버는 법",
        Inches(1), Inches(1.4), size=Pt(58))

# accent on "150만원"
rect(s, Inches(1), Inches(3.45), Inches(4.8), Inches(0.08), C_LIME)

txbox(s, "단기임대 수익 시스템 · 4주 현장강의",
      Inches(1), Inches(5.6), Inches(8), Inches(0.5),
      size=Pt(18), color=C_MUTED)

txbox(s, "© 2026 STAYCLASS",
      Inches(1), Inches(6.8), Inches(4), Inches(0.4),
      size=Pt(12), color=RGBColor(0x33, 0x33, 0x33))

# right decorative number
txbox(s, "150", Inches(8.5), Inches(1.5), Inches(5), Inches(4),
      size=Pt(220), bold=True, color=RGBColor(0x16, 0x16, 0x16))


# ─────────────────────────────────────────────
# SLIDE 2 — 목차
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_LIME)

label(s, "오늘 배울 것", Inches(1), Inches(0.7))
title_h(s, "4주 완성\n커리큘럼", Inches(1), Inches(1.1), size=Pt(48))

weeks = [
    ("1주차", "수익 나는 물건 찾기"),
    ("2주차", "계약 & 방 세팅"),
    ("3주차", "운영 자동화 구축"),
    ("4주차", "수익 극대화 & 확장"),
]

for i, (wk, ttl) in enumerate(weeks):
    cx = Inches(1 + i * 3.1)
    rect(s, cx, Inches(3.8), Inches(2.8), Inches(0.04), C_LIME)
    txbox(s, wk, cx, Inches(3.95), Inches(2.8), Inches(0.4),
          size=Pt(13), bold=True, color=C_LIME)
    txbox(s, ttl, cx, Inches(4.45), Inches(2.8), Inches(0.6),
          size=Pt(18), bold=True, color=C_TEXT)

txbox(s, "목표: 4주 후 방 1개 운영 시작 · 방 3개 → 월 150만원",
      Inches(1), Inches(6.5), Inches(10), Inches(0.5),
      size=Pt(16), color=C_MUTED)


# ─────────────────────────────────────────────
# SLIDE 3 — 왜 단기임대인가
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_LIME)

label(s, "왜 단기임대인가", Inches(1), Inches(0.7))
title_h(s, "이게 되는\n이유 셋", Inches(1), Inches(1.1), size=Pt(48))

reasons = [
    ("01", "내 집 없어도 됩니다",   "임차인 구조로 운영 · 초기비용 100~150만원"),
    ("02", "자동화가 됩니다",       "하루 30분 이내 운영 · 직장 병행 가능"),
    ("03", "월세보다 2배 납니다",    "월세 30~40만원 vs 단기임대 50~70만원"),
]

for i, (num, ttl, desc) in enumerate(reasons):
    cy = Inches(3.3 + i * 1.1)
    txbox(s, num, Inches(1), cy, Inches(0.7), Inches(0.8),
          size=Pt(14), bold=True, color=C_BORDER)
    txbox(s, ttl, Inches(1.7), cy, Inches(5), Inches(0.5),
          size=Pt(20), bold=True, color=C_TEXT)
    txbox(s, desc, Inches(1.7), cy + Inches(0.45), Inches(7), Inches(0.4),
          size=Pt(15), color=C_MUTED)
    if i < 2:
        rect(s, Inches(1), cy + Inches(0.95), Inches(11.33), Pt(0.5), C_BORDER)


# ─────────────────────────────────────────────
# SLIDE 4 — 수익 구조
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_LIME)

label(s, "수익 구조", Inches(1), Inches(0.7))
title_h(s, "방 3개 =\n월 150만원", Inches(1), Inches(1.1), size=Pt(48))

bars = [
    ("방 1개", "월 50~70만원", 0.38, "실측 기준", C_TEXT),
    ("방 3개", "월 150~210만원", 0.68, "목표", C_LIME),
    ("방 7개", "월 381만원", 1.0,  "강사 실측", C_BLUE),
]

bar_x = Inches(1)
bar_w_full = Inches(9)
for i, (lbl, val, pct, note, col) in enumerate(bars):
    cy = Inches(3.5 + i * 1.1)
    txbox(s, lbl, bar_x, cy, Inches(2), Inches(0.38),
          size=Pt(14), color=C_MUTED)
    txbox(s, val, Inches(3.2), cy, Inches(3.5), Inches(0.38),
          size=Pt(18), bold=True, color=col)
    txbox(s, note, Inches(6.8), cy + Inches(0.05), Inches(2), Inches(0.38),
          size=Pt(12), color=C_LIME if note == "실측 기준" else C_BLUE if note == "강사 실측" else C_MUTED)
    rect(s, bar_x, cy + Inches(0.48), bar_w_full, Inches(0.06), C_BORDER)
    rect(s, bar_x, cy + Inches(0.48), int(bar_w_full * pct), Inches(0.06), col)


# ─────────────────────────────────────────────
# SLIDE 5 — 실증 데이터
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_BLUE)

label(s, "강사 직접 운영 실증", Inches(1), Inches(0.7), color=C_BLUE)
title_h(s, "숫자가\n증명합니다", Inches(1), Inches(1.1), size=Pt(48))

# Data card background
rect(s, Inches(0.8), Inches(3.3), Inches(11.73), Inches(3.3),
     C_SURF, C_BORDER, Pt(1))
rect(s, Inches(0.8), Inches(3.3), Inches(11.73), Inches(0.05), C_BLUE)

txbox(s, "● LIVE  ·  Hostier 2026년 6월 기준", Inches(1), Inches(3.45),
      Inches(6), Inches(0.35), size=Pt(12), bold=True, color=C_BLUE)

stats = [
    ("90", "%", "가동률"),
    ("381", "만원", "월 순수익"),
    ("54", "만원", "방당 평균"),
    ("7", "개", "직영 방"),
]

for i, (num, unit, desc) in enumerate(stats):
    cx = Inches(1 + i * 3.0)
    txbox(s, num, cx, Inches(3.95), Inches(1.8), Inches(1.1),
          size=Pt(52), bold=True, color=C_LIME)
    txbox(s, unit, cx + Inches(1.2), Inches(4.1), Inches(1), Inches(0.5),
          size=Pt(18), color=C_BLUE)
    txbox(s, desc, cx, Inches(5.1), Inches(2.4), Inches(0.4),
          size=Pt(13), color=C_MUTED)
    if i < 3:
        rect(s, cx + Inches(2.7), Inches(4.0), Pt(1), Inches(1.2), C_BORDER)

txbox(s, "가동률 90% · 210박 중 188박 예약 확정",
      Inches(1), Inches(6.1), Inches(8), Inches(0.4),
      size=Pt(13), color=C_MUTED)


# ─────────────────────────────────────────────
# SLIDE 6 — 1주차
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_LIME)

label(s, "1주차", Inches(1), Inches(0.7))
title_h(s, "수익 나는\n물건 찾기", Inches(1), Inches(1.1), size=Pt(48))

items_1 = [
    "역세권 · 병원 · 법원 인근 지역 선정 공식",
    "삼삼엠투 vs 에어비앤비 수익 구조 비교",
    "수익 계산기로 사전 검증하는 법",
    "집주인 협상 스크립트 제공",
    "계약 전 체크리스트 완성",
]
bullet_list(s, items_1, Inches(1), Inches(3.3), Inches(8.5), size=Pt(19))

rect(s, Inches(9.8), Inches(1.3), Inches(2.9), Inches(4.5),
     C_SURF, C_BORDER, Pt(1))
txbox(s, "핵심 공식", Inches(9.9), Inches(1.5), Inches(2.6), Inches(0.4),
      size=Pt(12), bold=True, color=C_LIME)
txbox(s, "수요 = 유동인구\n÷ 숙소 수",
      Inches(9.9), Inches(2.0), Inches(2.6), Inches(1.2),
      size=Pt(20), bold=True, color=C_TEXT)
txbox(s, "이 비율이 높은 곳을\n먼저 선점합니다",
      Inches(9.9), Inches(3.3), Inches(2.6), Inches(0.8),
      size=Pt(13), color=C_MUTED)


# ─────────────────────────────────────────────
# SLIDE 7 — 2주차
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_LIME)

label(s, "2주차", Inches(1), Inches(0.7))
title_h(s, "계약 &\n방 세팅", Inches(1), Inches(1.1), size=Pt(48))

items_2 = [
    "계약서 조항 체크리스트 (전대차 허용 조항 필수)",
    "초기 세팅 물품 리스트 & 비용 최소화 전략",
    "사진 촬영으로 클릭률 높이는 법",
    "첫 예약 빠르게 받는 초기 가격 전략",
    "삼삼엠투 리스팅 첫 등록 실습",
]
bullet_list(s, items_2, Inches(1), Inches(3.3), Inches(8.5), size=Pt(19))

rect(s, Inches(9.8), Inches(1.3), Inches(2.9), Inches(3.8),
     C_SURF, C_BORDER, Pt(1))
txbox(s, "초기 비용 목표", Inches(9.9), Inches(1.5), Inches(2.6), Inches(0.4),
      size=Pt(12), bold=True, color=C_LIME)
costs = [("보증금", "50~100만원"), ("세팅 용품", "30~50만원"), ("중개비", "10~20만원")]
for i, (k, v) in enumerate(costs):
    txbox(s, k, Inches(9.9), Inches(2.1 + i * 0.65), Inches(1.3), Inches(0.5),
          size=Pt(13), color=C_MUTED)
    txbox(s, v, Inches(11.1), Inches(2.1 + i * 0.65), Inches(1.5), Inches(0.5),
          size=Pt(13), bold=True, color=C_TEXT)


# ─────────────────────────────────────────────
# SLIDE 8 — 3주차
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_LIME)

label(s, "3주차", Inches(1), Inches(0.7))
title_h(s, "운영\n자동화 구축", Inches(1), Inches(1.1), size=Pt(48))

items_3 = [
    "비대면 체크인 시스템 구축 (키박스 / 스마트락)",
    "청소 위탁 업체 선정 기준 & 단가 협상",
    "메시지 자동화 템플릿 전체 제공",
    "예약 확인 → 안내 → 체크아웃 자동 흐름 설계",
    "하루 30분 운영 루틴 만들기",
]
bullet_list(s, items_3, Inches(1), Inches(3.3), Inches(8.5), size=Pt(19))

rect(s, Inches(9.8), Inches(1.3), Inches(2.9), Inches(2.8),
     C_SURF, C_BORDER, Pt(1))
txbox(s, "자동화 후", Inches(9.9), Inches(1.5), Inches(2.6), Inches(0.4),
      size=Pt(12), bold=True, color=C_LIME)
txbox(s, "하루 30분", Inches(9.9), Inches(2.0), Inches(2.6), Inches(0.8),
      size=Pt(32), bold=True, color=C_LIME)
txbox(s, "직장 퇴근 후\n루틴 관리 가능",
      Inches(9.9), Inches(2.9), Inches(2.6), Inches(0.7),
      size=Pt(13), color=C_MUTED)


# ─────────────────────────────────────────────
# SLIDE 9 — 4주차
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_LIME)

label(s, "4주차", Inches(1), Inches(0.7))
title_h(s, "수익 극대화\n& 확장", Inches(1), Inches(1.1), size=Pt(48))

items_4 = [
    "시즌별 가격 조정 전략 (성수기 / 비수기)",
    "공실 최소화 — 다채널 & 할인 타이밍 전략",
    "플랫폼 노출 랭킹 올리는 법",
    "2호점 계약 타이밍 공식 (1호점 안정 후 60일)",
    "위탁운영 전환으로 완전 자동화 완성",
]
bullet_list(s, items_4, Inches(1), Inches(3.3), Inches(8.5), size=Pt(19))

rect(s, Inches(9.8), Inches(1.3), Inches(2.9), Inches(3.5),
     C_SURF, C_BORDER, Pt(1))
txbox(s, "최종 목표", Inches(9.9), Inches(1.5), Inches(2.6), Inches(0.4),
      size=Pt(12), bold=True, color=C_LIME)
txbox(s, "방 3개\n월 150만", Inches(9.9), Inches(2.0), Inches(2.6), Inches(1.2),
      size=Pt(32), bold=True, color=C_LIME)
txbox(s, "방 7개\n월 381만 (실증)",
      Inches(9.9), Inches(3.3), Inches(2.6), Inches(0.9),
      size=Pt(16), color=C_BLUE)


# ─────────────────────────────────────────────
# SLIDE 10 — Q&A
# ─────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, Inches(0.06), H, C_LIME)

txbox(s, "Q&A", Inches(1), Inches(1.5), Inches(10), Inches(2.5),
      size=Pt(110), bold=True, color=RGBColor(0x1A, 0x1A, 0x1A))

txbox(s, "궁금한 것 뭐든지 물어보세요",
      Inches(1), Inches(4.8), Inches(8), Inches(0.6),
      size=Pt(26), bold=True, color=C_TEXT)

txbox(s, "카카오채널 · pf.kakao.com/_APzxbX/chat",
      Inches(1), Inches(5.6), Inches(8), Inches(0.4),
      size=Pt(15), color=C_MUTED)

txbox(s, "STAYCLASS", Inches(1), Inches(6.6), Inches(4), Inches(0.4),
      size=Pt(13), bold=True, color=RGBColor(0x33, 0x33, 0x33))

# Lime accent bottom bar
rect(s, 0, H - Inches(0.08), W, Inches(0.08), C_LIME)


# ─────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────
out = r"C:\Users\User\test\stayclass_lecture.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
